"""
Office 문서(pptx, docx, xlsx)를 Markdown으로 변환하는 도구.
이미지를 추출하고 Groq VLM(Llama 4 Scout)으로 설명을 생성합니다.

사용법:
    python tools/convert_to_markdown.py [input_dir] [output_dir]

환경변수:
    GROQ_API_KEY: Groq API 키 (이미지 설명 생성에 필요)

기본값:
    input_dir:  resources/references
    output_dir: resources/references/markdown
"""

import sys
import os
import base64
import time
import io
from pathlib import Path
from dotenv import load_dotenv

# tools/.env 파일에서 환경변수 로드
load_dotenv(Path(__file__).resolve().parent / ".env")

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document
from docx.opc.constants import RELATIONSHIP_TYPE as RT
from openpyxl import load_workbook
from openpyxl.drawing.image import Image as XlImage

# ---------------------------------------------------------------------------
# Groq VLM
# ---------------------------------------------------------------------------

GROQ_MODEL = "meta-llama/llama-4-scout-17b-16e-instruct"
VLM_PROMPT = (
    "이미지에 포함된 모든 텍스트, 다이어그램, UI 요소, 캐릭터, 아이콘 등을 "
    "빠짐없이 상세하게 한국어로 설명해 주세요. "
    "텍스트가 있으면 그대로 옮겨 적고, 시각적 레이아웃과 구조도 설명해 주세요."
)

_groq_client = None


def _get_groq_client():
    global _groq_client
    if _groq_client is None:
        api_key = os.environ.get("GROQ_API_KEY")
        if not api_key:
            return None
        from groq import Groq
        _groq_client = Groq(api_key=api_key)
    return _groq_client


def describe_image(image_bytes: bytes, content_type: str = "image/png") -> str:
    """Groq VLM으로 이미지를 설명."""
    client = _get_groq_client()
    if client is None:
        return "(GROQ_API_KEY 미설정 - 이미지 설명 생략)"

    # 4MB 제한 체크
    if len(image_bytes) > 4 * 1024 * 1024:
        # 큰 이미지는 리사이즈
        image_bytes, content_type = _resize_image(image_bytes, max_bytes=3_500_000)

    b64 = base64.b64encode(image_bytes).decode("utf-8")
    data_url = f"data:{content_type};base64,{b64}"

    max_retries = 3
    for attempt in range(max_retries):
        try:
            response = client.chat.completions.create(
                model=GROQ_MODEL,
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": VLM_PROMPT},
                            {"type": "image_url", "image_url": {"url": data_url}},
                        ],
                    }
                ],
                temperature=0.3,
                max_completion_tokens=2048,
            )
            return response.choices[0].message.content.strip()
        except Exception as e:
            err_str = str(e).lower()
            if "rate_limit" in err_str or "429" in err_str:
                wait = 15 * (attempt + 1)
                print(f"\n    [Rate limit] {wait}초 대기 중...", end="", flush=True)
                time.sleep(wait)
            elif "413" in err_str or "too large" in err_str:
                image_bytes, content_type = _resize_image(image_bytes, max_bytes=2_000_000)
                b64 = base64.b64encode(image_bytes).decode("utf-8")
                data_url = f"data:{content_type};base64,{b64}"
            else:
                return f"(이미지 설명 실패: {e})"
    return "(이미지 설명 실패: rate limit 초과)"


MIN_IMAGE_WIDTH = 150   # 이 미만이면 아이콘으로 간주
MIN_IMAGE_HEIGHT = 150
MIN_IMAGE_BYTES = 5_000  # 5KB 미만이면 아이콘으로 간주


def _is_meaningful_image(image_bytes: bytes) -> bool:
    """아이콘/장식 이미지를 걸러내고 정보성 이미지만 통과."""
    if len(image_bytes) < MIN_IMAGE_BYTES:
        return False
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(image_bytes))
        w, h = img.size
        if w < MIN_IMAGE_WIDTH or h < MIN_IMAGE_HEIGHT:
            return False
    except Exception:
        pass
    return True


def _resize_image(image_bytes: bytes, max_bytes: int = 3_500_000) -> tuple[bytes, str]:
    """PIL로 이미지를 축소하여 크기 제한 내로 맞춤."""
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(image_bytes))
        quality = 85
        while True:
            buf = io.BytesIO()
            if img.mode in ("RGBA", "P"):
                img = img.convert("RGB")
            img.save(buf, format="JPEG", quality=quality)
            if buf.tell() <= max_bytes or quality <= 20:
                return buf.getvalue(), "image/jpeg"
            quality -= 15
            if quality <= 20:
                # 해상도도 줄이기
                w, h = img.size
                img = img.resize((w // 2, h // 2), Image.LANCZOS)
                quality = 70
    except Exception:
        return image_bytes, "image/png"


# ---------------------------------------------------------------------------
# 이미지 추출 헬퍼
# ---------------------------------------------------------------------------

def _save_image(image_bytes: bytes, ext: str, images_dir: Path, prefix: str, idx: int) -> Path:
    """이미지를 파일로 저장하고 경로 반환."""
    images_dir.mkdir(parents=True, exist_ok=True)
    if not ext.startswith("."):
        ext = "." + ext
    filename = f"{prefix}_img{idx:02d}{ext}"
    path = images_dir / filename
    path.write_bytes(image_bytes)
    return path


def _content_type_from_ext(ext: str) -> str:
    mapping = {
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".gif": "image/gif",
        ".bmp": "image/bmp",
        ".tiff": "image/tiff",
        ".tif": "image/tiff",
        ".wmf": "image/wmf",
        ".emf": "image/emf",
        ".svg": "image/svg+xml",
    }
    return mapping.get(ext.lower(), "image/png")


def _image_markdown(rel_path: str, description: str) -> list[str]:
    """이미지 참조와 설명을 마크다운 라인 리스트로 반환."""
    lines = [
        f"![이미지]({rel_path})\n",
        f"> {description}\n",
    ]
    return lines


# ---------------------------------------------------------------------------
# PPTX 변환
# ---------------------------------------------------------------------------

def _extract_shapes_images(shapes, images_dir: Path, prefix: str, img_counter: list, rel_base: str) -> list[str]:
    """셰이프 목록에서 이미지를 추출하고 마크다운 라인 반환."""
    img_lines = []
    for shape in shapes:
        # 그룹 내부 재귀
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            img_lines.extend(
                _extract_shapes_images(shape.shapes, images_dir, prefix, img_counter, rel_base)
            )
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE or (
            hasattr(shape, "image") and shape.shape_type not in (
                MSO_SHAPE_TYPE.TABLE, MSO_SHAPE_TYPE.CHART
            )
        ):
            try:
                img_blob = shape.image.blob
                img_ext = "." + shape.image.content_type.split("/")[-1]
                if img_ext == ".jpeg":
                    img_ext = ".jpg"

                img_counter[0] += 1
                saved = _save_image(img_blob, img_ext, images_dir, prefix, img_counter[0])
                rel_path = f"{rel_base}/{saved.name}"

                if not _is_meaningful_image(img_blob):
                    print(f"\n    [SKIP] {saved.name} (아이콘/장식)", flush=True)
                    img_lines.append(f"![아이콘]({rel_path})\n")
                    continue

                ct = _content_type_from_ext(img_ext)
                print(f"\n    [VLM] {saved.name} 분석 중...", end="", flush=True)
                desc = describe_image(img_blob, ct)
                print(" 완료", flush=True)

                img_lines.extend(_image_markdown(rel_path, desc))
            except Exception:
                pass
    return img_lines


def pptx_to_markdown(filepath: Path, output_dir: Path) -> str:
    """PowerPoint 파일을 Markdown으로 변환 (이미지 포함)."""
    prs = Presentation(str(filepath))
    doc_stem = filepath.stem
    images_dir = output_dir / "images" / doc_stem
    rel_base = f"images/{doc_stem}"

    lines = [f"# {doc_stem}\n"]
    img_counter = [0]

    for slide_idx, slide in enumerate(prs.slides, 1):
        lines.append(f"## 슬라이드 {slide_idx}\n")

        # 텍스트/테이블 먼저
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    rows.append(cells)

                if rows:
                    lines.append("| " + " | ".join(rows[0]) + " |")
                    lines.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
                    for row in rows[1:]:
                        while len(row) < len(rows[0]):
                            row.append("")
                        lines.append("| " + " | ".join(row[:len(rows[0])]) + " |")
                    lines.append("")

            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for child in shape.shapes:
                    if child.has_text_frame:
                        _extract_text_frame(child.text_frame, lines)

            elif shape.has_text_frame:
                _extract_text_frame(shape.text_frame, lines)

        # 이미지
        img_lines = _extract_shapes_images(
            slide.shapes, images_dir, f"slide{slide_idx:02d}", img_counter, rel_base
        )
        if img_lines:
            lines.append("")
            lines.extend(img_lines)

        lines.append("---\n")

    return "\n".join(lines)


def _extract_text_frame(text_frame, lines: list):
    """텍스트 프레임에서 마크다운 텍스트를 추출."""
    for para in text_frame.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        level = para.level if para.level else 0
        font_size = None
        is_bold = False
        if para.runs:
            font_size = para.runs[0].font.size
            is_bold = para.runs[0].font.bold or False

        if font_size and font_size >= Pt(24):
            lines.append(f"### {text}\n")
        elif font_size and font_size >= Pt(18):
            lines.append(f"#### {text}\n")
        elif is_bold:
            if level > 0:
                lines.append(f"{'  ' * level}- **{text}**")
            else:
                lines.append(f"**{text}**\n")
        else:
            if level > 0:
                lines.append(f"{'  ' * level}- {text}")
            else:
                lines.append(f"{text}\n")


# ---------------------------------------------------------------------------
# DOCX 변환
# ---------------------------------------------------------------------------

def docx_to_markdown(filepath: Path, output_dir: Path) -> str:
    """Word 문서를 Markdown으로 변환 (이미지 포함)."""
    doc = Document(str(filepath))
    doc_stem = filepath.stem
    images_dir = output_dir / "images" / doc_stem
    rel_base = f"images/{doc_stem}"

    lines = [f"# {doc_stem}\n"]

    # 이미지 추출 (relationships에서)
    img_counter = 0
    image_map = {}  # rId -> (rel_path, description)
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                img_blob = rel.target_part.blob
                ext = Path(rel.target_ref).suffix or ".png"
                img_counter += 1
                saved = _save_image(img_blob, ext, images_dir, "img", img_counter)
                rel_path = f"{rel_base}/{saved.name}"

                if not _is_meaningful_image(img_blob):
                    print(f"\n    [SKIP] {saved.name} (아이콘/장식)", flush=True)
                    image_map[rel.rId] = (rel_path, None)
                    continue

                ct = _content_type_from_ext(ext)
                print(f"\n    [VLM] {saved.name} 분석 중...", end="", flush=True)
                desc = describe_image(img_blob, ct)
                print(" 완료", flush=True)

                image_map[rel.rId] = (rel_path, desc)
            except Exception:
                pass

    # 문단 순서대로 처리 (이미지 위치 추적)
    for para in doc.paragraphs:
        text = para.text.strip()

        # 인라인 이미지 체크
        for run in para.runs:
            inline_shapes = run._element.findall(
                ".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing"
            )
            for drawing in inline_shapes:
                blip_elems = drawing.findall(
                    ".//{http://schemas.openxmlformats.org/drawingml/2006/main}blip"
                )
                for blip in blip_elems:
                    embed = blip.get(
                        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
                    )
                    if embed and embed in image_map:
                        rel_path, desc = image_map.pop(embed)
                        if desc is None:
                            lines.append(f"![아이콘]({rel_path})\n")
                        else:
                            lines.extend(_image_markdown(rel_path, desc))

        if not text:
            lines.append("")
            continue

        style_name = para.style.name.lower() if para.style else ""

        if "heading 1" in style_name or style_name == "제목 1":
            lines.append(f"## {text}\n")
        elif "heading 2" in style_name or style_name == "제목 2":
            lines.append(f"### {text}\n")
        elif "heading 3" in style_name or style_name == "제목 3":
            lines.append(f"#### {text}\n")
        elif "heading" in style_name:
            lines.append(f"##### {text}\n")
        elif "list" in style_name or "bullet" in style_name:
            lines.append(f"- {text}")
        else:
            formatted = _format_docx_paragraph(para)
            lines.append(f"{formatted}\n")

    # 본문에 매칭되지 않은 나머지 이미지 추가
    if image_map:
        lines.append("\n## 기타 이미지\n")
        for rel_path, desc in image_map.values():
            if desc is None:
                lines.append(f"![아이콘]({rel_path})\n")
            else:
                lines.extend(_image_markdown(rel_path, desc))

    # 테이블 처리
    for table in doc.tables:
        lines.append("")
        rows = []
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            rows.append(cells)

        if rows:
            col_count = len(rows[0])
            lines.append("| " + " | ".join(rows[0]) + " |")
            lines.append("| " + " | ".join(["---"] * col_count) + " |")
            for row in rows[1:]:
                while len(row) < col_count:
                    row.append("")
                lines.append("| " + " | ".join(row[:col_count]) + " |")
        lines.append("")

    return "\n".join(lines)


def _format_docx_paragraph(para) -> str:
    """Word 문단의 인라인 서식을 Markdown으로 변환."""
    parts = []
    for run in para.runs:
        text = run.text
        if not text:
            continue
        if run.bold and run.italic:
            parts.append(f"***{text}***")
        elif run.bold:
            parts.append(f"**{text}**")
        elif run.italic:
            parts.append(f"*{text}*")
        else:
            parts.append(text)
    return "".join(parts) if parts else para.text


# ---------------------------------------------------------------------------
# XLSX 변환
# ---------------------------------------------------------------------------

def xlsx_to_markdown(filepath: Path, output_dir: Path) -> str:
    """Excel 파일을 Markdown으로 변환 (이미지 포함)."""
    wb = load_workbook(str(filepath), data_only=True)
    doc_stem = filepath.stem
    images_dir = output_dir / "images" / doc_stem
    rel_base = f"images/{doc_stem}"

    lines = [f"# {doc_stem}\n"]
    img_counter = 0

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines.append(f"## {sheet_name}\n")

        # 시트 내 이미지
        if hasattr(ws, "_images") and ws._images:
            for xl_img in ws._images:
                try:
                    img_data = xl_img._data()
                    ext = Path(xl_img.path).suffix if hasattr(xl_img, "path") else ".png"
                    if not ext:
                        ext = ".png"
                    img_counter += 1
                    saved = _save_image(img_data, ext, images_dir, f"{sheet_name}_img", img_counter)
                    rel_path = f"{rel_base}/{saved.name}"

                    if not _is_meaningful_image(img_data):
                        print(f"\n    [SKIP] {saved.name} (아이콘/장식)", flush=True)
                        lines.append(f"![아이콘]({rel_path})\n")
                        continue

                    ct = _content_type_from_ext(ext)
                    print(f"\n    [VLM] {saved.name} 분석 중...", end="", flush=True)
                    desc = describe_image(img_data, ct)
                    print(" 완료", flush=True)

                    lines.extend(_image_markdown(rel_path, desc))
                except Exception:
                    pass

        # 테이블 데이터
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(cell) if cell is not None else "" for cell in row]
            if all(c == "" for c in cells):
                continue
            rows.append(cells)

        if not rows:
            lines.append("*(빈 시트)*\n")
            continue

        max_cols = max(len(r) for r in rows)
        for row in rows:
            while len(row) < max_cols:
                row.append("")

        lines.append("| " + " | ".join(rows[0]) + " |")
        lines.append("| " + " | ".join(["---"] * max_cols) + " |")
        for row in rows[1:]:
            escaped = [c.replace("|", "\\|") for c in row]
            lines.append("| " + " | ".join(escaped) + " |")
        lines.append("")

    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

CONVERTERS = {
    ".pptx": pptx_to_markdown,
    ".docx": docx_to_markdown,
    ".xlsx": xlsx_to_markdown,
}


def main():
    base_dir = Path(__file__).resolve().parent.parent
    input_dir = Path(sys.argv[1]) if len(sys.argv) > 1 else base_dir / "resources" / "references"
    output_dir = Path(sys.argv[2]) if len(sys.argv) > 2 else base_dir / "resources" / "references" / "markdown"

    output_dir.mkdir(parents=True, exist_ok=True)

    if not os.environ.get("GROQ_API_KEY"):
        print("경고: GROQ_API_KEY가 설정되지 않았습니다. 이미지 설명이 생략됩니다.\n")

    files = sorted(input_dir.iterdir())
    converted = 0
    errors = []
    total_images = 0

    for f in files:
        ext = f.suffix.lower()
        if ext not in CONVERTERS:
            continue

        print(f"변환 중: {f.name} ...", end=" ", flush=True)
        try:
            md_content = CONVERTERS[ext](f, output_dir)
            out_path = output_dir / (f.stem + ".md")
            out_path.write_text(md_content, encoding="utf-8")

            # 이미지 수 카운트
            img_count = md_content.count("![이미지]")
            total_images += img_count
            print(f"-> {out_path.name} (이미지 {img_count}개)")
            converted += 1
        except Exception as e:
            print(f"오류: {e}")
            errors.append((f.name, str(e)))

    print(f"\n완료: {converted}개 문서 변환, 총 이미지 {total_images}개 처리, {len(errors)}개 오류")
    if errors:
        print("\n오류 목록:")
        for name, err in errors:
            print(f"  - {name}: {err}")


if __name__ == "__main__":
    main()
